from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Color Palete ──────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
PANEL     = RGBColor(0x16, 0x2B, 0x3C)   # slightly lighter panel
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ORANGE    = RGBColor(0xFF, 0x9F, 0x1C)   # orange highlight
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xB0, 0xC4, 0xD8)   # light gray text
GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # green
YELLOW    = RGBColor(0xFF, 0xD1, 0x66)   # yellow

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, border_color=None,
             border_width=Pt(0), radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    shape.text = ""
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb


def add_label_box(slide, text, x, y, w, h, bg_color, txt_color=WHITE,
                  font_size=16, bold=True):
    """Colored rectangle with centered text."""
    add_rect(slide, x, y, w, h, fill_color=bg_color)
    add_textbox(slide, text, x, y + (h / 2) - 0.18, w, 0.4,
                font_size=font_size, bold=bold, color=txt_color,
                align=PP_ALIGN.CENTER)


def header_bar(slide, title, subtitle=""):
    """Top cyan accent bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=ACCENT)
    add_textbox(slide, title, 0.3, 0.1, 12.5, 0.7,
                font_size=32, bold=True, color=BG, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.72, 12.5, 0.4,
                    font_size=16, bold=False, color=BG, align=PP_ALIGN.LEFT)


def bullet_block(slide, items, x, y, w, icon="▶", icon_color=ACCENT,
                 font_size=17, line_gap=0.52):
    """Render a list of bullet strings."""
    for i, item in enumerate(items):
        yy = y + i * line_gap
        add_textbox(slide, icon, x, yy, 0.35, 0.45,
                    font_size=font_size, bold=True, color=icon_color)
        add_textbox(slide, item, x + 0.35, yy, w - 0.35, 0.45,
                    font_size=font_size, color=LGRAY)


def arrow(slide, x, y, length=0.7, horizontal=True):
    """Draw a simple line+arrowhead approximation using a textbox."""
    if horizontal:
        add_textbox(slide, "──►", x, y, length, 0.35,
                    font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, "▼", x, y, 0.35, length,
                    font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)

# gradient-like top strip
add_rect(sl, 0, 0, 13.33, 0.12, fill_color=ACCENT)
# bottom strip
add_rect(sl, 0, 7.38, 13.33, 0.12, fill_color=ORANGE)

# big title
add_textbox(sl, "Model Context Protocol", 1, 1.6, 11, 1.3,
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "(MCP)", 1, 2.85, 11, 0.8,
            font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# divider line — thin rect
add_rect(sl, 3.5, 3.7, 6.3, 0.05, fill_color=ORANGE)

add_textbox(sl, "What is MCP  ·  Why We Need It  ·  Architecture",
            1, 3.85, 11, 0.55,
            font_size=20, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

add_textbox(sl, "Study Guide", 1, 6.6, 11, 0.4,
            font_size=14, color=ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — What is MCP?
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "What is MCP?", "Model Context Protocol — a quick overview")

# big definition box
add_rect(sl, 0.4, 1.35, 12.5, 1.25, fill_color=PANEL,
         border_color=ACCENT, border_width=Pt(1.5))
add_textbox(sl,
    "MCP (Model Context Protocol) is an open standard created by Anthropic\n"
    "that lets AI models talk to external tools, files, databases & services.",
    0.6, 1.45, 12.1, 1.0, font_size=18, color=WHITE)

# three key idea cards
cards = [
    (ACCENT,  "📡  Open Standard",   "A universal protocol — like HTTP but for AI integrations."),
    (ORANGE,  "🔌  Plug & Play",      "Connect any AI app to any tool with one standard interface."),
    (GREEN,   "🤖  AI Superpower",    "Gives AI real-world awareness: files, APIs, databases & more."),
]
for i, (col, title, desc) in enumerate(cards):
    cx = 0.4 + i * 4.3
    add_rect(sl, cx, 2.8, 4.0, 2.4, fill_color=PANEL, border_color=col, border_width=Pt(2))
    add_textbox(sl, title, cx + 0.15, 2.95, 3.7, 0.5,
                font_size=17, bold=True, color=col)
    add_textbox(sl, desc, cx + 0.15, 3.5, 3.7, 1.5,
                font_size=15, color=LGRAY)

# analogy
add_rect(sl, 0.4, 5.35, 12.5, 0.9, fill_color=RGBColor(0x0F, 0x2E, 0x1E))
add_textbox(sl, "💡  Analogy: MCP is the USB-C port of AI — one standard connector "
               "for everything, instead of a different cable for every device.",
            0.6, 5.4, 12.1, 0.8, font_size=16, color=GREEN)

add_textbox(sl, "2", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — The Problem (Why we NEED MCP)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "The Problem Before MCP", "Why AI felt limited without it")

# Before panel
add_rect(sl, 0.3, 1.35, 5.9, 5.7, fill_color=PANEL,
         border_color=RGBColor(0xFF, 0x5C, 0x5C), border_width=Pt(1.5))
add_textbox(sl, "❌  Before MCP", 0.5, 1.45, 5.5, 0.5,
            font_size=20, bold=True, color=RGBColor(0xFF, 0x5C, 0x5C))

before_items = [
    "AI models were isolated — no access to files or internet",
    "Every tool needed its own custom integration code",
    "No shared standard → each team reinvented the wheel",
    "Hard to maintain & scale across many tools",
    "AI answers were limited to training data only",
]
bullet_block(sl, before_items, 0.5, 2.05, 5.5,
             icon="✗", icon_color=RGBColor(0xFF, 0x5C, 0x5C), font_size=15)

# After panel
add_rect(sl, 7.1, 1.35, 5.9, 5.7, fill_color=PANEL,
         border_color=GREEN, border_width=Pt(1.5))
add_textbox(sl, "✅  With MCP", 7.3, 1.45, 5.5, 0.5,
            font_size=20, bold=True, color=GREEN)

after_items = [
    "AI connects to any tool via one standard protocol",
    "Build a server once → reuse across any AI app",
    "Consistent interface for tools, files & databases",
    "Easy to add new capabilities without rewriting code",
    "AI becomes truly context-aware & useful",
]
bullet_block(sl, after_items, 7.3, 2.05, 5.5,
             icon="✓", icon_color=GREEN, font_size=15)

# center VS
add_textbox(sl, "VS", 6.1, 3.8, 1.1, 0.6,
            font_size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_textbox(sl, "3", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — Need for MCP (6 reasons)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "Why Do We Need MCP?", "6 key reasons MCP matters")

reasons = [
    (ACCENT,  "1. Standardisation",      "One protocol for all AI ↔ tool communication"),
    (ORANGE,  "2. Reusability",           "Write a server once, use it with any AI host"),
    (GREEN,   "3. Real-world Access",     "Let AI read files, query DBs, call APIs"),
    (YELLOW,  "4. Less Code",             "No more custom glue code for every integration"),
    (RGBColor(0xC7, 0x7D, 0xFF), "5. Scalability", "Easily add new tools without breaking existing ones"),
    (RGBColor(0xFF, 0x6B, 0x9D), "6. Ecosystem",   "Community servers → thousands of ready-made tools"),
]

for i, (col, title, desc) in enumerate(reasons):
    col_idx = i % 3
    row_idx = i // 3
    cx = 0.35 + col_idx * 4.3
    cy = 1.45 + row_idx * 2.6
    add_rect(sl, cx, cy, 4.1, 2.3, fill_color=PANEL, border_color=col, border_width=Pt(2))
    # top color strip
    add_rect(sl, cx, cy, 4.1, 0.12, fill_color=col)
    add_textbox(sl, title, cx + 0.12, cy + 0.2, 3.8, 0.5,
                font_size=16, bold=True, color=col)
    add_textbox(sl, desc, cx + 0.12, cy + 0.75, 3.8, 1.3,
                font_size=14, color=LGRAY)

add_textbox(sl, "4", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — Architecture Overview (diagram)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "MCP Architecture", "The three-layer model")

# ── Row 1: HOST
add_rect(sl, 0.3, 1.4, 12.7, 1.8, fill_color=RGBColor(0x0A, 0x25, 0x38),
         border_color=ACCENT, border_width=Pt(1))
add_textbox(sl, "HOST  (e.g. Claude Desktop, VS Code, Cursor)",
            0.45, 1.45, 6, 0.45, font_size=13, bold=True, color=ACCENT)

# Client boxes inside host
for i, label in enumerate(["MCP Client 1", "MCP Client 2", "MCP Client 3"]):
    cx = 0.5 + i * 4.1
    add_rect(sl, cx, 1.95, 3.8, 0.95, fill_color=RGBColor(0x1A, 0x3A, 0x52),
             border_color=ACCENT, border_width=Pt(1))
    add_textbox(sl, label, cx, 2.1, 3.8, 0.6,
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Transport label
add_textbox(sl, "◄──  Transport Layer  ──►   (stdio  |  HTTP + SSE)",
            0.3, 3.3, 12.7, 0.4,
            font_size=14, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

# ── Row 2: SERVERS
server_data = [
    (ACCENT,  "MCP Server A",  "File System\nTool"),
    (GREEN,   "MCP Server B",  "Database\nTool"),
    (ORANGE,  "MCP Server C",  "Web Search\nTool"),
]
for i, (col, name, sub) in enumerate(server_data):
    cx = 0.5 + i * 4.1
    add_rect(sl, cx, 3.75, 3.8, 1.5, fill_color=PANEL,
             border_color=col, border_width=Pt(2))
    add_rect(sl, cx, 3.75, 3.8, 0.12, fill_color=col)
    add_textbox(sl, name, cx, 3.92, 3.8, 0.45,
                font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(sl, sub, cx, 4.42, 3.8, 0.75,
                font_size=13, color=LGRAY, align=PP_ALIGN.CENTER)

# ── Row 3: Resources / Tools / Prompts
labels = ["📄  Resources\n(Files, DB rows)", "🔧  Tools\n(Run functions)", "📝  Prompts\n(Templates)"]
cols_r = [ACCENT, GREEN, ORANGE]
for i, (lbl, col) in enumerate(zip(labels, cols_r)):
    cx = 0.5 + i * 4.1
    add_rect(sl, cx, 5.4, 3.8, 1.0, fill_color=RGBColor(0x0A, 0x1F, 0x2E),
             border_color=col, border_width=Pt(1))
    add_textbox(sl, lbl, cx, 5.52, 3.8, 0.8,
                font_size=14, color=col, align=PP_ALIGN.CENTER)

add_textbox(sl, "5", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — MCP Components Deep Dive
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "MCP Components — Deep Dive", "What each part does")

components = [
    (ACCENT,  "🖥  Host",
     "The AI application that the user interacts with.\n"
     "Examples: Claude Desktop, VS Code with Copilot, Cursor IDE.\n"
     "It contains one or more MCP Clients."),
    (ORANGE,  "🔗  Client",
     "Lives inside the Host. Manages the connection to an MCP Server.\n"
     "Sends requests (tool calls, resource reads) and receives responses.\n"
     "One client ↔ one server."),
    (GREEN,   "⚙  Server",
     "A lightweight program that exposes capabilities to the AI.\n"
     "Can provide Tools (actions), Resources (data), or Prompts (templates).\n"
     "You build this to connect AI to your custom system."),
    (YELLOW,  "📡  Transport",
     "How Client and Server communicate:\n"
     "• stdio — local process communication (simple)\n"
     "• HTTP + SSE — network communication (remote servers)"),
]

for i, (col, title, body) in enumerate(components):
    row = i // 2
    col_idx = i % 2
    cx = 0.35 + col_idx * 6.5
    cy = 1.45 + row * 2.9
    add_rect(sl, cx, cy, 6.2, 2.65, fill_color=PANEL,
             border_color=col, border_width=Pt(2))
    add_rect(sl, cx, cy, 6.2, 0.12, fill_color=col)
    add_textbox(sl, title, cx + 0.15, cy + 0.18, 5.8, 0.5,
                font_size=18, bold=True, color=col)
    add_textbox(sl, body, cx + 0.15, cy + 0.72, 5.8, 1.8,
                font_size=14, color=LGRAY)

add_textbox(sl, "6", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — How MCP Works (Step-by-step flow)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "How MCP Works", "Step-by-step communication flow")

steps = [
    (ACCENT,  "1",  "User asks AI",          "User types a question in the Host app\n(e.g. \"Summarise my notes.txt\")"),
    (ORANGE,  "2",  "AI decides",            "AI figures out it needs an external tool\nor resource to answer the question"),
    (GREEN,   "3",  "Client calls Server",   "MCP Client sends a request to the\nappropriate MCP Server"),
    (YELLOW,  "4",  "Server responds",       "Server fetches the file / runs the function\nand sends the result back"),
    (RGBColor(0xC7,0x7D,0xFF), "5", "AI answers", "AI uses the result to give a\nfull, accurate answer to the user"),
]

box_w = 2.3
box_h = 3.2
start_x = 0.35
y_top = 1.45

for i, (col, num, title, desc) in enumerate(steps):
    cx = start_x + i * (box_w + 0.15)
    # connector arrow between boxes
    if i > 0:
        add_textbox(sl, "►", cx - 0.17, y_top + box_h / 2 - 0.2, 0.2, 0.4,
                    font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(sl, cx, y_top, box_w, box_h, fill_color=PANEL,
             border_color=col, border_width=Pt(2))
    # number circle (just a square here)
    add_rect(sl, cx + box_w/2 - 0.32, y_top + 0.15, 0.64, 0.64, fill_color=col)
    add_textbox(sl, num, cx + box_w/2 - 0.32, y_top + 0.18, 0.64, 0.45,
                font_size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, cx + 0.1, y_top + 0.92, box_w - 0.2, 0.5,
                font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, cx + 0.1, y_top + 1.5, box_w - 0.2, 1.55,
                font_size=12, color=LGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, "7", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — Key Concepts (Tools / Resources / Prompts)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "Key MCP Concepts", "The three things an MCP Server can expose")

concepts = [
    (ACCENT,  "🔧  Tools",
     "Functions the AI can CALL to do something.",
     ["Search the web", "Send an email", "Write to a file",
      "Run a database query", "Call an external API"]),
    (GREEN,   "📄  Resources",
     "Data the AI can READ (like files or DB rows).",
     ["Read a text file", "Get a webpage", "Fetch DB records",
      "Read config settings", "Access logs"]),
    (ORANGE,  "📝  Prompts",
     "Pre-built prompt templates for common tasks.",
     ["Code review template", "Summary prompt",
      "Bug report format", "Meeting notes template"]),
]

for i, (col, title, subtitle, examples) in enumerate(concepts):
    cx = 0.35 + i * 4.3
    add_rect(sl, cx, 1.4, 4.1, 5.6, fill_color=PANEL,
             border_color=col, border_width=Pt(2))
    add_rect(sl, cx, 1.4, 4.1, 0.12, fill_color=col)
    add_textbox(sl, title, cx + 0.15, 1.55, 3.8, 0.55,
                font_size=20, bold=True, color=col)
    add_textbox(sl, subtitle, cx + 0.15, 2.15, 3.8, 0.55,
                font_size=14, color=LGRAY, italic=True)
    add_rect(sl, cx + 0.15, 2.75, 3.8, 0.04, fill_color=col)
    add_textbox(sl, "Examples:", cx + 0.15, 2.85, 3.8, 0.4,
                font_size=13, bold=True, color=col)
    bullet_block(sl, examples, cx + 0.1, 3.3, 3.9,
                 icon="•", icon_color=col, font_size=13, line_gap=0.45)

add_textbox(sl, "8", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — Summary & Key Takeaways
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
header_bar(sl, "Summary", "What you should remember")

takeaways = [
    (ACCENT,  "MCP = Standard bridge",      "Connects AI models to external tools, files & services"),
    (ORANGE,  "3 Core Components",          "Host  →  Client  →  Server"),
    (GREEN,   "3 Capabilities",             "Tools (do)  ·  Resources (read)  ·  Prompts (template)"),
    (YELLOW,  "2 Transport Options",        "stdio (local)  and  HTTP+SSE (remote)"),
    (RGBColor(0xFF,0x6B,0x9D), "Why it matters", "One protocol replaces hundreds of custom integrations"),
    (RGBColor(0xC7,0x7D,0xFF), "Analogy",        "MCP is the USB-C port of AI applications"),
]

for i, (col, title, desc) in enumerate(takeaways):
    row = i // 2
    col_idx = i % 2
    cx = 0.35 + col_idx * 6.5
    cy = 1.5 + row * 1.75
    add_rect(sl, cx, cy, 6.2, 1.5, fill_color=PANEL,
             border_color=col, border_width=Pt(2))
    add_rect(sl, cx, cy, 0.12, 1.5, fill_color=col)  # left accent stripe
    add_textbox(sl, title, cx + 0.25, cy + 0.12, 5.8, 0.45,
                font_size=16, bold=True, color=col)
    add_textbox(sl, desc, cx + 0.25, cy + 0.6, 5.8, 0.75,
                font_size=14, color=LGRAY)

add_textbox(sl, "9", 12.9, 7.1, 0.35, 0.3, font_size=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════════
out_path = r"D:\Study\MCP\MCP_Introduction.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
