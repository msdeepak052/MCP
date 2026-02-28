# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)
PANEL   = RGBColor(0x16, 0x2B, 0x3C)
PANEL2  = RGBColor(0x0F, 0x22, 0x30)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
ORANGE  = RGBColor(0xFF, 0x9F, 0x1C)
GREEN   = RGBColor(0x06, 0xD6, 0xA0)
YELLOW  = RGBColor(0xFF, 0xD1, 0x66)
PURPLE  = RGBColor(0xC7, 0x7D, 0xFF)
PINK    = RGBColor(0xFF, 0x6B, 0x9D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xB0, 0xC4, 0xD8)
RED     = RGBColor(0xFF, 0x5C, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    s.text = ""
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    r.font.name      = name
    return tb

def header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.15, fill=ACCENT)
    txt(slide, title, 0.3, 0.08, 12.5, 0.65, size=32, bold=True, color=BG)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.70, 12.5, 0.38, size=15, color=BG, italic=True)

def slide_num(slide, n):
    txt(slide, str(n), 12.85, 7.1, 0.4, 0.3, size=12, color=LGRAY, align=PP_ALIGN.RIGHT)

def bullets(slide, items, x, y, w, icon=">>", icol=ACCENT, size=15, gap=0.50):
    for i, item in enumerate(items):
        yy = y + i * gap
        txt(slide, icon, x, yy, 0.38, 0.42, size=size, bold=True, color=icol)
        txt(slide, item, x + 0.38, yy, w - 0.38, 0.42, size=size, color=LGRAY)

def card(slide, x, y, w, h, col, title, body, title_size=16, body_size=14):
    rect(slide, x, y, w, h, fill=PANEL, border=col, bw=Pt(2))
    rect(slide, x, y, w, 0.10, fill=col)
    txt(slide, title, x+0.15, y+0.15, w-0.3, 0.48, size=title_size, bold=True, color=col)
    txt(slide, body,  x+0.15, y+0.68, w-0.3, h-0.8, size=body_size, color=LGRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
rect(sl, 0, 0,     13.33, 0.10, fill=ACCENT)
rect(sl, 0, 7.40,  13.33, 0.10, fill=ORANGE)

txt(sl, "What is MCP?", 0.5, 1.3, 12.3, 1.4, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Model Context Protocol — A Complete Guide",
    0.5, 2.75, 12.3, 0.65, size=24, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
rect(sl, 3.8, 3.55, 5.7, 0.06, fill=ORANGE)
txt(sl, "Architecture  |  Components  |  Transport  |  Lifecycle  |  Use Cases",
    0.5, 3.72, 12.3, 0.5, size=17, color=LGRAY, align=PP_ALIGN.CENTER)
txt(sl, "Comprehensive Study Guide", 0.5, 6.7, 12.3, 0.4,
    size=14, color=ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — MCP in One Line
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "MCP in One Line", "The simplest explanation")

rect(sl, 0.4, 1.3, 12.5, 1.4, fill=PANEL2, border=ACCENT, bw=Pt(2))
txt(sl,
    '"MCP is an open standard that lets AI applications connect to external\n'
    ' tools, data sources, and services in a structured, secure way."',
    0.65, 1.38, 12.0, 1.22, size=20, color=WHITE, italic=True)

# Three analogy cards
analogy_data = [
    (ACCENT,  "USB-C Port",
     "Just like USB-C is one universal\nconnector for all devices,\nMCP is one protocol for all AI tools."),
    (ORANGE,  "HTTP for AI",
     "HTTP made the web work by\nstandarsing communication.\nMCP does the same for AI tools."),
    (GREEN,   "Plugin System",
     "Like browser extensions add\ncapabilities to your browser,\nMCP servers add powers to AI."),
]
for i, (col, title, desc) in enumerate(analogy_data):
    cx = 0.4 + i * 4.3
    rect(sl, cx, 2.9, 4.1, 3.9, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, 2.9, 4.1, 0.10, fill=col)
    txt(sl, "Analogy:", cx+0.15, 3.05, 3.8, 0.36, size=12, color=LGRAY, italic=True)
    txt(sl, title, cx+0.15, 3.42, 3.8, 0.5, size=20, bold=True, color=col)
    txt(sl, desc, cx+0.15, 4.0, 3.8, 2.5, size=14, color=LGRAY)
slide_num(sl, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Origin
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Background & Origin", "Where MCP came from")

# Timeline-style horizontal layout
timeline = [
    (ACCENT,  "Nov 2024",      "Anthropic releases MCP as an open standard to the public."),
    (ORANGE,  "The Problem",   "AI models couldn't connect to tools without custom code for each integration."),
    (GREEN,   "The Vision",    "One universal protocol so any AI app can work with any tool."),
    (YELLOW,  "Open Source",   "MCP spec & SDKs are fully open source — anyone can build servers."),
    (PURPLE,  "Adoption",      "Quickly adopted by Cursor, Zed, Replit, Codeium and many others."),
]
for i, (col, label, desc) in enumerate(timeline):
    cy = 1.35 + i * 1.15
    rect(sl, 0.35, cy, 2.2, 0.85, fill=col)
    txt(sl, label, 0.35, cy+0.18, 2.2, 0.5, size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(sl, 2.55, cy+0.35, 0.5, 0.06, fill=col)   # connector
    rect(sl, 3.05, cy, 9.9, 0.85, fill=PANEL, border=col, bw=Pt(1))
    txt(sl, desc, 3.2, cy+0.18, 9.5, 0.52, size=14, color=LGRAY)
slide_num(sl, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — The Problem MCP Solves
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "The Problem MCP Solves", "Why AI needed a standard protocol")

# Top: the chaos diagram
rect(sl, 0.35, 1.3, 12.6, 2.6, fill=PANEL2, border=RED, bw=Pt(1))
txt(sl, "BEFORE MCP — Every tool needed its own custom glue code",
    0.55, 1.35, 12.0, 0.45, size=15, bold=True, color=RED)

# AI box in center
rect(sl, 5.5, 1.95, 2.3, 0.9, fill=RGBColor(0x1A,0x1A,0x2E), border=RED, bw=Pt(2))
txt(sl, "AI Model", 5.5, 2.12, 2.3, 0.55, size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

tools_before = ["File System", "Database", "Web API", "Email", "Slack", "GitHub"]
cols_b = [0.45, 1.9, 3.35, 8.0, 9.45, 10.9]
for j, (tool, cx) in enumerate(zip(tools_before, cols_b)):
    rect(sl, cx, 2.0, 1.3, 0.8, fill=PANEL, border=LGRAY, bw=Pt(1))
    txt(sl, tool, cx, 2.18, 1.3, 0.45, size=11, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(sl, "custom\ncode", cx+0.2, 2.55, 0.9, 0.35, size=9, color=RED, align=PP_ALIGN.CENTER, italic=True)

txt(sl, "Result: N tools = N custom integrations. Hard to build, maintain & scale.",
    0.55, 3.58, 12.0, 0.4, size=13, color=RED, italic=True)

# Bottom: MCP solution
rect(sl, 0.35, 4.1, 12.6, 2.85, fill=PANEL2, border=GREEN, bw=Pt(1))
txt(sl, "WITH MCP — One standard protocol for everything",
    0.55, 4.15, 12.0, 0.45, size=15, bold=True, color=GREEN)

rect(sl, 5.5, 4.7, 2.3, 0.9, fill=RGBColor(0x0A,0x25,0x15), border=GREEN, bw=Pt(2))
txt(sl, "AI Model", 5.5, 4.88, 2.3, 0.55, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txt(sl, "MCP Protocol", 4.5, 5.72, 4.3, 0.38, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
rect(sl, 4.5, 5.75, 4.3, 0.06, fill=ACCENT)

tools_after = ["File Srv", "DB Srv", "Web Srv", "Email Srv", "Slack Srv", "Git Srv"]
for j, (tool, cx) in enumerate(zip(tools_after, cols_b)):
    rect(sl, cx, 6.0, 1.3, 0.7, fill=PANEL, border=GREEN, bw=Pt(1))
    txt(sl, tool, cx, 6.17, 1.3, 0.38, size=11, color=GREEN, align=PP_ALIGN.CENTER)
slide_num(sl, 4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture Overview
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "MCP Architecture Overview", "The big picture")

# HOST container
rect(sl, 0.3, 1.25, 12.73, 2.5, fill=RGBColor(0x08,0x1C,0x2B), border=ACCENT, bw=Pt(1.5))
txt(sl, "HOST APPLICATION  (e.g. Claude Desktop / Cursor / VS Code)",
    0.45, 1.3, 8, 0.42, size=13, bold=True, color=ACCENT)

# Clients inside host
for i, lbl in enumerate(["MCP Client 1", "MCP Client 2", "MCP Client 3"]):
    cx = 0.5 + i * 4.1
    rect(sl, cx, 1.82, 3.8, 1.6, fill=RGBColor(0x12,0x2B,0x3E), border=ACCENT, bw=Pt(1))
    txt(sl, lbl, cx, 2.28, 3.8, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "manages connection", cx, 2.78, 3.8, 0.42,
        size=11, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

# Transport band
rect(sl, 0.3, 3.85, 12.73, 0.55, fill=RGBColor(0x1A,0x0F,0x00), border=ORANGE, bw=Pt(1))
txt(sl, "TRANSPORT LAYER  --  stdio (local)   |   HTTP + SSE (remote)",
    0.5, 3.9, 12.3, 0.38, size=14, color=ORANGE, align=PP_ALIGN.CENTER, bold=True)

# Server boxes
server_info = [
    (ACCENT, "MCP Server A", "File System"),
    (GREEN,  "MCP Server B", "Database"),
    (ORANGE, "MCP Server C", "Web / APIs"),
]
for i, (col, name, sub) in enumerate(server_info):
    cx = 0.5 + i * 4.1
    rect(sl, cx, 4.55, 3.8, 1.35, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, 4.55, 3.8, 0.10, fill=col)
    txt(sl, name, cx, 4.72, 3.8, 0.48, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, sub,  cx, 5.22, 3.8, 0.42, size=13, color=LGRAY, align=PP_ALIGN.CENTER)

# Capabilities row
cap_data = [
    (ACCENT, "Tools"), (GREEN, "Resources"), (ORANGE, "Prompts"),
    (ACCENT, "Tools"), (GREEN, "Resources"), (ORANGE, "Prompts"),
    (ACCENT, "Tools"), (GREEN, "Resources"), (ORANGE, "Prompts"),
]
positions = []
for srv in range(3):
    for cap in range(3):
        positions.append((0.5 + srv*4.1 + cap*1.27, 6.05))

for k, ((col, lbl), (cx, cy)) in enumerate(zip(cap_data, positions)):
    rect(sl, cx, cy, 1.2, 0.55, fill=RGBColor(0x0A,0x20,0x2E), border=col, bw=Pt(1))
    txt(sl, lbl, cx, cy+0.1, 1.2, 0.38, size=10, color=col, align=PP_ALIGN.CENTER)
slide_num(sl, 5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — The Three Core Components
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "The Three Core Components", "Host, Client & Server explained")

comp_data = [
    (ACCENT, "HOST", "The AI Application",
     "The app the user directly uses.\n\n"
     "Examples: Claude Desktop, Cursor, VS Code\n\n"
     "It CONTAINS one or more MCP Clients.\n\n"
     "Responsibility:\n"
     "  - Manage the user interaction\n"
     "  - Hold & run MCP Clients\n"
     "  - Enforce security & consent"),
    (GREEN, "CLIENT", "The Bridge",
     "Lives INSIDE the Host.\n\n"
     "One Client <-> One Server (1:1 mapping)\n\n"
     "Responsibility:\n"
     "  - Open/maintain connection to a Server\n"
     "  - Send requests (tool calls, reads)\n"
     "  - Receive & relay responses to the AI\n"
     "  - Handle protocol negotiation"),
    (ORANGE, "SERVER", "The Capability Provider",
     "A separate program YOU build or install.\n\n"
     "Exposes: Tools, Resources, Prompts\n\n"
     "Examples: filesystem server, DB server\n\n"
     "Responsibility:\n"
     "  - Execute tool calls safely\n"
     "  - Serve resource data on request\n"
     "  - Can be local or remote"),
]
for i, (col, name, subtitle, body) in enumerate(comp_data):
    cx = 0.35 + i * 4.32
    rect(sl, cx, 1.25, 4.1, 5.8, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, 1.25, 4.1, 0.12, fill=col)
    txt(sl, name, cx+0.15, 1.38, 3.8, 0.52, size=22, bold=True, color=col)
    txt(sl, subtitle, cx+0.15, 1.9, 3.8, 0.4, size=14, color=LGRAY, italic=True)
    rect(sl, cx+0.15, 2.35, 3.8, 0.04, fill=col)
    txt(sl, body, cx+0.15, 2.45, 3.8, 4.4, size=13, color=LGRAY)
slide_num(sl, 6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Transport Layer
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Transport Layer", "How Client and Server communicate")

txt(sl, "Transport = the communication channel between MCP Client and MCP Server",
    0.4, 1.25, 12.5, 0.45, size=16, color=LGRAY, italic=True)

# stdio card
rect(sl, 0.35, 1.85, 6.0, 5.25, fill=PANEL, border=ACCENT, bw=Pt(2))
rect(sl, 0.35, 1.85, 6.0, 0.12, fill=ACCENT)
txt(sl, "stdio  (Standard Input/Output)", 0.5, 2.0, 5.7, 0.52, size=20, bold=True, color=ACCENT)
txt(sl, "Best for: LOCAL servers on same machine", 0.5, 2.55, 5.7, 0.38, size=13, color=GREEN, italic=True)

stdio_pts = [
    "Client launches Server as a child process",
    "They talk through stdin & stdout pipes",
    "Simple, fast, no network needed",
    "Server lives & dies with the Client",
    "Great for filesystem, local DBs, scripts",
]
bullets(sl, stdio_pts, 0.5, 3.05, 5.7, icon="->", icol=ACCENT, size=14, gap=0.48)

# diagram for stdio
rect(sl, 0.6, 5.75, 2.0, 0.6, fill=RGBColor(0x0A,0x25,0x38), border=ACCENT, bw=Pt(1))
txt(sl, "Client", 0.6, 5.88, 2.0, 0.35, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "---stdin/stdout--->", 2.65, 5.9, 2.5, 0.35, size=11, color=LGRAY, align=PP_ALIGN.CENTER)
rect(sl, 5.2, 5.75, 1.0, 0.6, fill=RGBColor(0x0A,0x25,0x38), border=ACCENT, bw=Pt(1))
txt(sl, "Srv", 5.2, 5.88, 1.0, 0.35, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# HTTP+SSE card
rect(sl, 6.98, 1.85, 6.0, 5.25, fill=PANEL, border=ORANGE, bw=Pt(2))
rect(sl, 6.98, 1.85, 6.0, 0.12, fill=ORANGE)
txt(sl, "HTTP + SSE", 7.13, 2.0, 5.7, 0.52, size=20, bold=True, color=ORANGE)
txt(sl, "Best for: REMOTE servers over a network", 7.13, 2.55, 5.7, 0.38, size=13, color=GREEN, italic=True)

http_pts = [
    "Client sends requests via HTTP POST",
    "Server pushes updates via SSE (Server-Sent Events)",
    "Server runs independently (not a child process)",
    "Can be shared by multiple Clients",
    "Great for cloud services, shared tools, SaaS",
]
bullets(sl, http_pts, 7.13, 3.05, 5.7, icon="->", icol=ORANGE, size=14, gap=0.48)

# diagram for HTTP
rect(sl, 7.15, 5.75, 2.0, 0.6, fill=RGBColor(0x25,0x18,0x00), border=ORANGE, bw=Pt(1))
txt(sl, "Client", 7.15, 5.88, 2.0, 0.35, size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, "--HTTP/SSE-->", 9.2, 5.9, 2.0, 0.35, size=11, color=LGRAY, align=PP_ALIGN.CENTER)
rect(sl, 11.25, 5.75, 1.6, 0.6, fill=RGBColor(0x25,0x18,0x00), border=ORANGE, bw=Pt(1))
txt(sl, "Remote Srv", 11.25, 5.88, 1.6, 0.35, size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
slide_num(sl, 7)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Three Capabilities (Tools / Resources / Prompts)
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "What an MCP Server Can Expose", "Three types of capabilities")

cap_data = [
    (ACCENT, "TOOLS", "AI can CALL these",
     "Functions that DO something.\nThe AI can invoke them to take action.",
     ["search_web(query)", "read_file(path)", "send_email(to, body)",
      "run_query(sql)", "create_issue(title)"],
     "Model-controlled -- AI decides when to call"),
    (GREEN, "RESOURCES", "AI can READ these",
     "Data that the AI can access.\nLike files or database rows.",
     ["file://notes.txt", "db://users/table", "screen://current",
      "git://repo/log", "http://api/data"],
     "App-controlled -- host decides what to expose"),
    (ORANGE, "PROMPTS", "Pre-built templates",
     "Reusable prompt templates for\ncommon workflows.",
     ["code_review template", "summarise_doc template",
      "debug_error template", "meeting_notes template",
      "write_tests template"],
     "User-controlled -- user selects the template"),
]

for i, (col, title, tag, desc, examples, note) in enumerate(cap_data):
    cx = 0.35 + i * 4.32
    rect(sl, cx, 1.25, 4.1, 5.85, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, 1.25, 4.1, 0.12, fill=col)
    txt(sl, title, cx+0.15, 1.38, 3.8, 0.52, size=22, bold=True, color=col)
    rect(sl, cx+0.15, 1.9, 1.8, 0.32, fill=col)
    txt(sl, tag, cx+0.16, 1.91, 1.78, 0.3, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, desc, cx+0.15, 2.3, 3.8, 0.75, size=13, color=LGRAY)
    txt(sl, "Examples:", cx+0.15, 3.12, 3.8, 0.35, size=13, bold=True, color=col)
    bullets(sl, examples, cx+0.1, 3.5, 3.95, icon="*", icol=col, size=12, gap=0.42)
    rect(sl, cx+0.15, 6.55, 3.8, 0.04, fill=col)
    txt(sl, note, cx+0.15, 6.62, 3.8, 0.38, size=11, color=LGRAY, italic=True)
slide_num(sl, 8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Tools Deep Dive
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Tools — Deep Dive", "The most powerful MCP capability")

rect(sl, 0.35, 1.28, 12.6, 0.82, fill=PANEL2, border=ACCENT, bw=Pt(1))
txt(sl, "A Tool is a function exposed by an MCP Server that the AI model can CALL to perform "
        "an action. The AI sends a tool call, the server executes it and returns the result.",
    0.55, 1.33, 12.2, 0.72, size=15, color=WHITE)

# Left: how tools work
rect(sl, 0.35, 2.25, 5.9, 4.85, fill=PANEL, border=ACCENT, bw=Pt(1.5))
txt(sl, "How Tools Work", 0.5, 2.32, 5.6, 0.46, size=17, bold=True, color=ACCENT)
steps = [
    ("1. Discover", "Server tells Client which tools are available + their parameters"),
    ("2. AI decides", "AI reads the tool list and picks the right tool for the task"),
    ("3. Call",      "AI sends: tool name + arguments (structured JSON)"),
    ("4. Execute",   "Server runs the function safely"),
    ("5. Result",    "Server returns the output back to the AI"),
    ("6. Answer",    "AI uses the result to respond to the user"),
]
for i, (step, desc) in enumerate(steps):
    cy = 2.85 + i * 0.72
    rect(sl, 0.5, cy, 1.1, 0.52, fill=ACCENT)
    txt(sl, step, 0.5, cy+0.08, 1.1, 0.36, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, desc, 1.7, cy+0.08, 4.4, 0.36, size=12, color=LGRAY)

# Right: example tool definition
rect(sl, 6.55, 2.25, 6.45, 4.85, fill=RGBColor(0x0A,0x14,0x0A), border=GREEN, bw=Pt(1.5))
txt(sl, "Example Tool Definition", 6.7, 2.32, 6.1, 0.46, size=17, bold=True, color=GREEN)
code = (
    'Tool: read_file\n'
    'Description:\n'
    '  Read the contents of a file\n\n'
    'Parameters:\n'
    '  path (string, required)\n'
    '    -- The file path to read\n\n'
    'Returns:\n'
    '  content (string)\n'
    '    -- File contents as text\n\n'
    'AI Call Example:\n'
    '  read_file(path="notes.txt")\n\n'
    'Server Returns:\n'
    '  "Meeting at 3pm. Buy milk."'
)
txt(sl, code, 6.7, 2.82, 6.1, 4.1, size=12, color=GREEN, name="Consolas")
slide_num(sl, 9)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Resources Deep Dive
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Resources — Deep Dive", "Data the AI can read")

rect(sl, 0.35, 1.28, 12.6, 0.82, fill=PANEL2, border=GREEN, bw=Pt(1))
txt(sl, "A Resource is a piece of data exposed by an MCP Server. The AI (or Host app) can "
        "READ it to gain context. Unlike Tools, Resources don't execute actions — they supply information.",
    0.55, 1.33, 12.2, 0.72, size=15, color=WHITE)

# Resource types grid
res_types = [
    (ACCENT,  "Text Resources",    "Plain text files, logs,\nconfig files, notes.\n\nURI: file:///path/to/file.txt"),
    (GREEN,   "Binary Resources",  "Images, PDFs, audio.\nReturned as base64.\n\nURI: file:///image.png"),
    (ORANGE,  "Dynamic Resources", "Real-time data like\nlive DB rows or API output.\n\nURI: db://table/rows"),
    (PURPLE,  "Resource Templates","Parameterised URIs\ne.g. file:///{filename}\nto access any matching file"),
]
for i, (col, title, desc) in enumerate(res_types):
    row, col_idx = i // 2, i % 2
    cx = 0.35 + col_idx * 6.2
    cy = 2.28 + row * 2.2
    rect(sl, cx, cy, 6.0, 2.0, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, cy, 6.0, 0.10, fill=col)
    txt(sl, title, cx+0.15, cy+0.16, 5.7, 0.46, size=17, bold=True, color=col)
    txt(sl, desc, cx+0.15, cy+0.68, 5.7, 1.2, size=14, color=LGRAY)

# bottom note
rect(sl, 0.35, 6.65, 12.6, 0.5, fill=RGBColor(0x0A,0x25,0x15), border=GREEN, bw=Pt(1))
txt(sl, "Resources use URI addressing  --  each resource has a unique URI like a web URL.",
    0.55, 6.7, 12.2, 0.38, size=14, color=GREEN)
slide_num(sl, 10)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Prompts Deep Dive
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Prompts — Deep Dive", "Reusable AI workflow templates")

rect(sl, 0.35, 1.28, 12.6, 0.82, fill=PANEL2, border=ORANGE, bw=Pt(1))
txt(sl, "Prompts are pre-built, reusable prompt templates that an MCP Server exposes. "
        "The USER selects a prompt, the Host fills in parameters, and sends it to the AI — "
        "giving consistent, repeatable workflows.",
    0.55, 1.33, 12.2, 0.72, size=15, color=WHITE)

# Left column
rect(sl, 0.35, 2.25, 5.9, 4.9, fill=PANEL, border=ORANGE, bw=Pt(1.5))
txt(sl, "Why Prompts Matter", 0.5, 2.32, 5.6, 0.46, size=17, bold=True, color=ORANGE)
why_pts = [
    "Consistent outputs every time",
    "Encode best practices in a template",
    "Users don't need to write complex prompts",
    "Can include context from Resources",
    "Teams can share proven workflows",
    "Discoverable -- AI can list available prompts",
]
bullets(sl, why_pts, 0.5, 2.9, 5.6, icon=">>", icol=ORANGE, size=14, gap=0.56)

# Right column: example
rect(sl, 6.55, 2.25, 6.45, 4.9, fill=RGBColor(0x1A,0x10,0x00), border=ORANGE, bw=Pt(1.5))
txt(sl, "Example Prompt Template", 6.7, 2.32, 6.1, 0.46, size=17, bold=True, color=ORANGE)
prompt_ex = (
    'Name: code_review\n'
    'Description:\n'
    '  Review code for bugs & improvements\n\n'
    'Parameters:\n'
    '  language (string)\n'
    '    -- Programming language\n'
    '  code (string)\n'
    '    -- The code to review\n\n'
    'Generated Prompt:\n'
    '  "Review this {language} code for\n'
    '   bugs, style issues, and suggest\n'
    '   improvements:\n'
    '   {code}"'
)
txt(sl, prompt_ex, 6.7, 2.82, 6.1, 4.1, size=12, color=ORANGE, name="Consolas")
slide_num(sl, 11)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — MCP Lifecycle
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "MCP Session Lifecycle", "What happens from start to finish")

phases = [
    (ACCENT,  "1. Initialize",
     "Client connects to Server.\n"
     "They exchange protocol versions\n"
     "and capabilities.\n"
     "Handshake is completed.",
     "Like a handshake before\na conversation begins."),
    (GREEN,   "2. Discover",
     "Client asks Server:\n"
     "- What tools do you have?\n"
     "- What resources can I read?\n"
     "- What prompts exist?\n"
     "Server lists all capabilities.",
     "Like reading a menu before\nordering food."),
    (ORANGE,  "3. Operate",
     "Normal working phase.\n"
     "AI calls tools as needed.\n"
     "Reads resources on demand.\n"
     "Server executes & returns\nresults. Repeats as needed.",
     "The actual work gets done\nduring this phase."),
    (YELLOW,  "4. Shutdown",
     "Client sends close signal.\n"
     "Server cleans up resources.\n"
     "Connection is closed cleanly.\n"
     "No data loss or corruption.",
     "Clean exit like closing\nan app properly."),
]

for i, (col, phase, desc, analogy) in enumerate(phases):
    cx = 0.35 + i * 3.22
    # connector arrow
    if i > 0:
        txt(sl, "==>", cx-0.22, 3.3, 0.25, 0.38, size=13, color=ACCENT, align=PP_ALIGN.CENTER)
    # main card
    rect(sl, cx, 1.28, 3.0, 4.0, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, 1.28, 3.0, 0.12, fill=col)
    txt(sl, phase, cx+0.12, 1.42, 2.76, 0.5, size=17, bold=True, color=col)
    txt(sl, desc,  cx+0.12, 2.0,  2.76, 2.6, size=13, color=LGRAY)
    # analogy strip
    rect(sl, cx+0.12, 4.5, 2.76, 0.04, fill=col)
    txt(sl, "Analogy:", cx+0.12, 4.6, 2.76, 0.32, size=11, bold=True, color=col)
    txt(sl, analogy, cx+0.12, 4.95, 2.76, 0.72, size=12, color=LGRAY, italic=True)

# bottom note
rect(sl, 0.35, 5.8, 12.6, 0.62, fill=PANEL2, border=LGRAY, bw=Pt(1))
txt(sl, "Note: MCP uses JSON-RPC 2.0 as the message format for all communication between Client and Server.",
    0.55, 5.88, 12.2, 0.46, size=14, color=LGRAY, italic=True)
slide_num(sl, 12)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — MCP vs Function Calling
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "MCP vs Function Calling", "What's the difference?")

# Table header
rect(sl, 0.35, 1.28, 4.5, 0.55, fill=LGRAY)
rect(sl, 4.85, 1.28, 3.9, 0.55, fill=PURPLE)
rect(sl, 8.75, 1.28, 4.22, 0.55, fill=ACCENT)
txt(sl, "Aspect", 0.5, 1.33, 4.2, 0.4, size=15, bold=True, color=BG)
txt(sl, "Function Calling", 5.0, 1.33, 3.6, 0.4, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "MCP (Tools)", 8.9, 1.33, 3.9, 0.4, size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)

rows = [
    ("What is it?",         "Built into the LLM API.\nAI calls a function directly.",          "External protocol.\nAI calls a server-hosted function."),
    ("Where does it live?", "Inside your app code.\nFunctions defined per-app.",               "In a separate MCP Server.\nReusable across any app."),
    ("Reusability",         "Low -- tied to one app.\nMust rewrite for each project.",          "High -- one server works with\nany MCP-compatible AI app."),
    ("Standardisation",     "No standard -- each LLM\nprovider has their own format.",         "Open standard -- same protocol\nfor all AI models & apps."),
    ("Best for",            "Simple, one-off integrations\nwithin a single app.",               "Complex, reusable tools shared\nacross multiple AI applications."),
]
row_cols = [PANEL, PANEL2, PANEL, PANEL2, PANEL]
for i, (aspect, fc, mcp) in enumerate(rows):
    cy = 1.9 + i * 1.07
    rc = row_cols[i]
    rect(sl, 0.35, cy, 4.5, 1.0, fill=rc)
    rect(sl, 4.85, cy, 3.9, 1.0, fill=rc)
    rect(sl, 8.75, cy, 4.22, 1.0, fill=rc)
    # borders
    rect(sl, 0.35, cy, 4.5, 1.0, border=LGRAY, bw=Pt(0.5))
    rect(sl, 4.85, cy, 3.9, 1.0, border=PURPLE, bw=Pt(0.5))
    rect(sl, 8.75, cy, 4.22, 1.0, border=ACCENT, bw=Pt(0.5))
    txt(sl, aspect, 0.5, cy+0.08, 4.2, 0.85, size=13, bold=True, color=WHITE)
    txt(sl, fc,     5.0, cy+0.08, 3.6, 0.85, size=12, color=LGRAY)
    txt(sl, mcp,    8.9, cy+0.08, 3.9, 0.85, size=12, color=LGRAY)
slide_num(sl, 13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Real-World Use Cases
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Real-World Use Cases", "What can you actually DO with MCP?")

use_cases = [
    (ACCENT,  "Developer Tools",
     "- AI reads your codebase files\n"
     "- Runs tests & shows results\n"
     "- Creates GitHub issues/PRs\n"
     "- Searches documentation"),
    (GREEN,   "Data & Databases",
     "- AI queries your SQL database\n"
     "- Reads CSV or JSON files\n"
     "- Generates reports from data\n"
     "- Updates records safely"),
    (ORANGE,  "Communication",
     "- AI drafts & sends emails\n"
     "- Posts to Slack channels\n"
     "- Creates calendar events\n"
     "- Reads meeting notes"),
    (PURPLE,  "Web & Research",
     "- AI searches the web live\n"
     "- Fetches webpage content\n"
     "- Scrapes data from sites\n"
     "- Monitors news / prices"),
    (YELLOW,  "File Management",
     "- Reads, writes, moves files\n"
     "- Organises folder structure\n"
     "- Converts file formats\n"
     "- Creates backups"),
    (PINK,    "Custom Business",
     "- Connect to your internal API\n"
     "- Access proprietary data\n"
     "- Automate workflows\n"
     "- Integrate legacy systems"),
]
for i, (col, title, desc) in enumerate(use_cases):
    row, col_idx = i // 3, i % 3
    cx = 0.35 + col_idx * 4.32
    cy = 1.28 + row * 2.9
    rect(sl, cx, cy, 4.1, 2.7, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, cy, 4.1, 0.10, fill=col)
    txt(sl, title, cx+0.15, cy+0.16, 3.8, 0.5, size=17, bold=True, color=col)
    txt(sl, desc,  cx+0.15, cy+0.72, 3.8, 1.8, size=13, color=LGRAY)
slide_num(sl, 14)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Popular MCP Servers (Ecosystem)
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "MCP Ecosystem", "Popular ready-made MCP Servers")

txt(sl, "Anthropic & the community have built hundreds of MCP servers you can use immediately:",
    0.4, 1.25, 12.5, 0.42, size=15, color=LGRAY, italic=True)

servers = [
    (ACCENT,  "Official by Anthropic"),
    (GREEN,   "Official by Anthropic"),
    (ORANGE,  "Official by Anthropic"),
    (ACCENT,  "Community"),
    (GREEN,   "Community"),
    (ORANGE,  "Community"),
    (PURPLE,  "Community"),
    (YELLOW,  "Community"),
    (PINK,    "Community"),
]
server_names = [
    ("Filesystem",   "Read/write local files\nsafely with permissions"),
    ("PostgreSQL",   "Query a Postgres DB\nwith read-only access"),
    ("GitHub",       "Manage repos, issues,\nPRs via GitHub API"),
    ("Brave Search", "Search the web using\nBrave Search API"),
    ("Slack",        "Post messages, read\nchannels, manage threads"),
    ("Puppeteer",    "Control a browser:\nscreenshots, scraping"),
    ("Notion",       "Read/write Notion\npages and databases"),
    ("Google Drive", "Access Google Docs,\nSheets, files"),
    ("SQLite",       "Query a local SQLite\ndatabase file"),
]

for i, ((col, badge), (name, desc)) in enumerate(zip(servers, server_names)):
    row, col_idx = i // 3, i % 3
    cx = 0.35 + col_idx * 4.32
    cy = 1.82 + row * 1.82
    rect(sl, cx, cy, 4.1, 1.65, fill=PANEL, border=col, bw=Pt(1.5))
    rect(sl, cx, cy, 4.1, 0.10, fill=col)
    txt(sl, name, cx+0.15, cy+0.16, 2.5, 0.46, size=16, bold=True, color=col)
    rect(sl, cx+2.8, cy+0.2, 1.15, 0.3, fill=col)
    txt(sl, badge, cx+2.8, cy+0.2, 1.15, 0.3, size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, desc, cx+0.15, cy+0.68, 3.8, 0.85, size=13, color=LGRAY)
slide_num(sl, 15)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Security in MCP
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Security in MCP", "How MCP stays safe")

# Left: principles
rect(sl, 0.35, 1.28, 6.1, 5.85, fill=PANEL, border=GREEN, bw=Pt(1.5))
txt(sl, "Security Principles", 0.5, 1.35, 5.8, 0.46, size=18, bold=True, color=GREEN)

sec_pts = [
    ("User Consent",     "Every tool call requires explicit user approval.\nThe user controls what the AI can do."),
    ("Least Privilege",  "Servers only get the permissions they need.\nFilesystem server can't access the network."),
    ("Sandboxing",       "MCP Servers run in isolation.\nOne server can't access another server."),
    ("Audit Trail",      "All tool calls are logged.\nYou can see exactly what the AI did."),
    ("No Secrets",       "MCP Servers should not store API keys\nin responses — keep secrets server-side."),
]
for i, (title, desc) in enumerate(sec_pts):
    cy = 1.92 + i * 0.98
    rect(sl, 0.5, cy, 0.5, 0.75, fill=GREEN)
    txt(sl, str(i+1), 0.5, cy+0.16, 0.5, 0.42, size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, title, 1.12, cy+0.05, 5.1, 0.35, size=14, bold=True, color=GREEN)
    txt(sl, desc,  1.12, cy+0.4,  5.1, 0.45, size=12, color=LGRAY)

# Right: trust model
rect(sl, 6.75, 1.28, 6.22, 5.85, fill=PANEL, border=ORANGE, bw=Pt(1.5))
txt(sl, "Trust Model", 6.9, 1.35, 5.9, 0.46, size=18, bold=True, color=ORANGE)

trust_data = [
    (ACCENT,  "HOST",   "Trusted completely.\nYou chose to install it."),
    (GREEN,   "CLIENT", "Trusted — it's inside\nthe Host you control."),
    (ORANGE,  "SERVER", "Partially trusted.\nUser approves each action."),
    (RED,     "User Input", "Prompt injection risk.\nValidate all inputs."),
]
for i, (col, name, note) in enumerate(trust_data):
    cy = 1.92 + i * 1.25
    rect(sl, 6.9, cy, 1.5, 1.05, fill=col)
    txt(sl, name, 6.9, cy+0.28, 1.5, 0.5, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, "Trust Level:", 8.55, cy+0.08, 4.2, 0.32, size=12, bold=True, color=col)
    txt(sl, note, 8.55, cy+0.45, 4.2, 0.5, size=12, color=LGRAY)
slide_num(sl, 16)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17 — Summary & Key Takeaways
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
header(sl, "Summary — Everything About MCP", "What you must remember")

takeaways = [
    (ACCENT,  "What",        "Open standard by Anthropic (Nov 2024) for AI-to-tool communication"),
    (GREEN,   "Why",         "Solves the N*M integration problem -- one protocol replaces custom code"),
    (ORANGE,  "Architecture","Host > Client > Server. Transport: stdio (local) or HTTP+SSE (remote)"),
    (YELLOW,  "Capabilities","Tools (do), Resources (read), Prompts (templates)"),
    (PURPLE,  "Lifecycle",   "Initialize > Discover > Operate > Shutdown. Uses JSON-RPC 2.0"),
    (PINK,    "vs Fn Calling","MCP is reusable & standard. Function calling is app-specific."),
    (ACCENT,  "Security",    "User consent, least privilege, sandboxing, audit trail"),
    (GREEN,   "Ecosystem",   "Hundreds of ready servers: GitHub, Postgres, Slack, Filesystem..."),
]

for i, (col, label, desc) in enumerate(takeaways):
    row, col_idx = i // 2, i % 2
    cx = 0.35 + col_idx * 6.5
    cy = 1.32 + row * 1.46
    rect(sl, cx, cy, 6.2, 1.3, fill=PANEL, border=col, bw=Pt(2))
    rect(sl, cx, cy, 0.12, 1.3, fill=col)
    rect(sl, cx+0.22, cy+0.12, 0.9, 0.36, fill=col)
    txt(sl, label, cx+0.22, cy+0.12, 0.9, 0.36, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, desc, cx+0.22, cy+0.55, 5.8, 0.65, size=14, color=LGRAY)

# footer
rect(sl, 0.35, 7.05, 12.6, 0.32, fill=ACCENT)
txt(sl, "MCP = The universal language between AI and the world.",
    0.5, 7.07, 12.2, 0.28, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
slide_num(sl, 17)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out = r"D:\Study\MCP\MCP_Complete_Guide.pptx"
prs.save(out)
print("Saved:", out)
